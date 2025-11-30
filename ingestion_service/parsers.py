import io
from pathlib import Path
from typing import Union

import pypdf
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

def parse_pdf(file_input: Union[str, Path, bytes]) -> str:
    """Extract text from a PDF file path or bytes."""
    # Convert bytes to stream if necessary
    if isinstance(file_input, bytes):
        stream = io.BytesIO(file_input)
        reader = pypdf.PdfReader(stream)
    else:
        # It's a path
        reader = pypdf.PdfReader(file_input)
        
    text = []
    for page in reader.pages:
        content = page.extract_text()
        if content:
            text.append(content)
    return "\n".join(text)

def parse_docx(file_input: Union[str, Path, bytes]) -> str:
    """Extract text from a .docx file path or bytes."""
    if isinstance(file_input, bytes):
        file_input = io.BytesIO(file_input)
        
    doc = Document(file_input)
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

def parse_pptx(file_input: Union[str, Path, bytes]) -> str:
    """Extract text from a .pptx file path or bytes."""
    if isinstance(file_input, bytes):
        file_input = io.BytesIO(file_input)
        
    prs = Presentation(file_input)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text)
    return "\n".join(text)

def parse_xlsx(file_input: Union[str, Path, bytes]) -> str:
    """Extract text from a .xlsx file path or bytes."""
    if isinstance(file_input, bytes):
        file_input = io.BytesIO(file_input)
        
    wb = load_workbook(file_input, read_only=True, data_only=True)
    text = []
    for sheet in wb.worksheets:
        sheet_text = []
        for row in sheet.iter_rows(values_only=True):
            # Convert row cells to string and join them
            row_str = " | ".join([str(cell) for cell in row if cell is not None])
            if row_str.strip():
                sheet_text.append(row_str)
        
        if sheet_text:
            text.append(f"--- Sheet: {sheet.title} ---")
            text.extend(sheet_text)
            
    return "\n".join(text)