import os
from pathlib import Path
from reportlab.pdfgen import canvas
from openpyxl import Workbook
from pptx import Presentation
from docx import Document

# Ensure the directory exists
DOCS_DIR = Path("data/docs")
DOCS_DIR.mkdir(parents=True, exist_ok=True)

def create_mongodb_pdf():
    filepath = DOCS_DIR / "mongodb_docs.pdf"
    print(f"Generating {filepath}...")
    
    c = canvas.Canvas(str(filepath))
    c.drawString(100, 800, "MongoDB Documentation Overview")
    c.drawString(100, 780, "--------------------------------")
    
    text_lines = [
        "MongoDB is a source-available cross-platform document-oriented database program.",
        "Classified as a NoSQL database program, MongoDB uses JSON-like documents with optional schemas.",
        "Key Features:",
        "1. Ad-hoc queries for search by field, range query, and regular expression searches.",
        "2. Indexing: Any field in a MongoDB document can be indexed.",
        "3. Replication: MongoDB provides high availability with replica sets.",
        "4. Load balancing: Sharding distributes data across a cluster of machines.",
    ]
    
    y = 750
    for line in text_lines:
        c.drawString(100, y, line)
        y -= 20
        
    c.save()

def create_nodejs_excel():
    filepath = DOCS_DIR / "nodejs_stats.xlsx"
    print(f"Generating {filepath}...")
    
    wb = Workbook()
    ws = wb.active
    ws.title = "NodeJS Stats"
    
    # Header
    ws.append(["Year", "Version", "Downloads (Millions)", "Active Developers (Millions)"])
    
    # Data rows
    data = [
        [2020, "14.x", 500, 8.5],
        [2021, "16.x", 750, 9.8],
        [2022, "18.x", 980, 11.2],
        [2023, "20.x", 1200, 13.5],
        [2024, "22.x", 1500, 15.0],
    ]
    
    for row in data:
        ws.append(row)
        
    wb.save(filepath)

def create_pydantic_pptx():
    filepath = DOCS_DIR / "pydantic_docs.pptx"
    print(f"Generating {filepath}...")
    
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Pydantic Documentation"
    subtitle.text = "Data Validation using Python Type Hints"
    
    # Slide 2: Content
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Why Pydantic?"
    
    content = slide.placeholders[1]
    content.text = (
        "1. Type Hints: Uses standard Python type hints.\n"
        "2. Speed: Core validation logic is written in Rust (pydantic-core).\n"
        "3. JSON Schema: Pydantic models can emit JSON Schema automatically.\n"
        "4. Ecosystem: Heavily used by FastAPI, LangChain, and other modern frameworks."
    )
    
    prs.save(filepath)

def create_word_doc():
    # Bonus: Creating a Word doc just to test the .docx parser too
    filepath = DOCS_DIR / "fastapi_extra.docx"
    print(f"Generating {filepath}...")
    
    doc = Document()
    doc.add_heading('FastAPI Extra Features', 0)
    doc.add_paragraph('FastAPI supports automatic interactive API documentation.')
    doc.add_paragraph('It is based on open standards: OpenAPI (previously known as Swagger) and JSON Schema.')
    doc.save(filepath)

if __name__ == "__main__":
    create_mongodb_pdf()
    create_nodejs_excel()
    create_pydantic_pptx()
    create_word_doc()
    print("\nAll test files generated in data/docs/")